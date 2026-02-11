from pptx import Presentation

def build_ppt(template_path, slide_json, output):
    prs = Presentation(template_path)

    # TITLE SLIDE
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = slide_json["title_slide"]["title"]

    # NORMAL SLIDES
    for s in slide_json["slides"]:
        layout_id = LAYOUT_MAP.get(s["type"], 1)
        slide_layout = prs.slide_layouts[layout_id]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = s["title"]

        if s["type"] == "bullet":
            tf = slide.placeholders[1].text_frame
            for p in s["points"]:
                tf.add_paragraph().text = p

    prs.save(output)